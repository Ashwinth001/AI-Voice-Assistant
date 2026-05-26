"""
Agent Dispatcher v2 - handles all actions.
No API keys needed for any feature.
"""
import os, sys, subprocess, webbrowser, threading, time, json, re, schedule
from pathlib import Path
from datetime import datetime
from core.config_loader import load_config

_cfg     = load_config()
CODE_DIR = Path(_cfg["agentic"]["code_output_dir"])
CODE_DIR.mkdir(parents=True, exist_ok=True)
APP_NAME = _cfg["app"]["name"]


class AgentDispatcher:
    def __init__(self, tts, llm):
        self.tts    = tts
        self.llm    = llm
        self._sched = []
        self._start_scheduler()

        # Lazy imports (optional integrations)
        self._win   = None
        self._learn = None
        self._screen= None

    def _win_integration(self):
        if self._win is None:
            try:
                from agents.windows_integration import OutlookIntegration, TeamsIntegration, WindowsAutomation
                self._outlook  = OutlookIntegration()
                self._teams    = TeamsIntegration()
                self._winauto  = WindowsAutomation()
                self._win = True
            except Exception as e:
                print(f"[Dispatcher] Windows integration not available: {e}")
                self._win = False
        return self._win

    def _get_learner(self):
        if self._learn is None:
            from core.web_learner import WebLearner
            self._learn = WebLearner(self.llm.memory)
        return self._learn

    def _get_screen(self):
        if self._screen is None:
            from core.screen_analyzer import analyze_screen
            self._screen = analyze_screen
        return self._screen

    def dispatch(self, intent: dict, lang: str = "en") -> bool:
        action = intent.get("action", "")
        params = intent.get("params", {})
        fn_map = {
            "open_browser":         self._browser,
            "search_web":           self._search,
            "open_file":            self._open_file,
            "open_folder":          self._open_folder,
            "create_file":          self._create_file,
            "create_folder":        self._create_folder,
            "write_code":           self._write_code,
            "run_code":             self._run_code,
            "set_reminder":         self._reminder,
            "create_ppt":           self._create_ppt,
            "research_topic":       self._research,
            "analyze_file":         self._analyze_file,
            "system_info":          self._sysinfo,
            "restart_pc":           self._restart,
            "shutdown_pc":          self._shutdown,
            "create_teams_meeting": self._teams_meeting,
            "check_email":          self._check_email,
            "send_email":           self._send_email,
            "screen_analyze":       self._screen_analyze,
            "update_model":         self._update_model,
            "open_notepad_dictate": self._dictate,
            # New commands
            "identity_response":    self._identity,
            "take_photo":           self._take_photo,
            "start_video":          self._start_video,
            "stop_video":           self._stop_video,
            "take_screenshot":      self._screenshot,
            "start_screen_recording": self._start_screen_rec,
            "stop_screen_recording":  self._stop_screen_rec,
            "auto_learn":           self._auto_learn,
        }
        fn = fn_map.get(action)
        if fn:
            threading.Thread(target=fn, args=(params, lang), daemon=True).start()
            return True
        return False

    #  Browser 
    def _browser(self, p, lang):
        url = p.get("url", "https://www.google.com")
        if not url.startswith("http"):
            url = "https://" + url
        self.tts.speak("Opening that now.", lang)
        webbrowser.open(url)

    def _search(self, p, lang):
        q = p.get("query", "")
        self.tts.speak(f"Searching for {q}.", lang)
        webbrowser.open(f"https://www.google.com/searchq={q.replace(' ','+')}")

    #  File system 
    def _open_file(self, p, lang):
        path = Path(p.get("path", ""))
        if path.exists():
            self.tts.speak("Opening file.", lang)
            os.startfile(str(path)) if sys.platform=="win32" else subprocess.run(["xdg-open",str(path)])
        else:
            self.tts.speak("File not found. Please check the path.", lang)

    def _open_folder(self, p, lang):
        path = p.get("path", str(Path.home()))
        self.tts.speak("Opening folder.", lang)
        os.startfile(path) if sys.platform=="win32" else subprocess.run(["xdg-open",path])

    def _create_file(self, p, lang):
        path = Path(p.get("path","new_file.txt"))
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(p.get("content",""), encoding="utf-8")
        self.tts.speak(f"Created {path.name}.", lang)

    def _create_folder(self, p, lang):
        path = Path(p.get("path","new_folder"))
        path.mkdir(parents=True, exist_ok=True)
        self.tts.speak(f"Folder created.", lang)

    #  Code 
    def _ext(self, l):
        return {"python":"py","go":"go","javascript":"js","java":"java",
                "rust":"rs","typescript":"ts","cpp":"cpp","c":"c"}.get(l,"py")

    def _write_code(self, p, lang):
        task     = p.get("task", "")
        clang    = p.get("language", "python")
        filename = p.get("filename", f"solution.{self._ext(clang)}")
        self.tts.speak(f"Writing {clang} code. Give me a moment.", lang)

        prompt = (f"Write complete working {clang} code for: {task}\n"
                  f"Return ONLY the raw code. No explanation. No markdown fences.")
        code = self.llm.simple_query(prompt)

        # Strip markdown fences
        code = re.sub(r"```[a-z]*\n", "", code).replace("```", "").strip()

        fpath = CODE_DIR / filename
        fpath.write_text(code, encoding="utf-8")
        self.llm.memory.store_code(code, clang, task)

        try:
            subprocess.Popen([_cfg["agentic"]["editor"], str(fpath)])
        except Exception:
            os.startfile(str(fpath)) if sys.platform=="win32" else None

        self.tts.speak(f"Done. Saved as {filename} and opened in your editor.", lang)

    def _run_code(self, p, lang):
        path  = p.get("path","")
        clang = p.get("language","python")
        cmds  = {"python":["python",path],"go":["go","run",path],"node":["node",path]}
        cmd   = cmds.get(clang, ["python",path])
        self.tts.speak("Running it now.", lang)
        try:
            r   = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
            out = (r.stdout or r.stderr or "No output.").strip()[:200]
            self.tts.speak(f"Result: {out}", lang)
        except Exception as e:
            self.tts.speak(f"Error: {str(e)[:80]}", lang)

    #  Reminder 
    def _reminder(self, p, lang):
        msg     = p.get("message","reminder")
        minutes = int(p.get("minutes", 5))
        def fire():
            self.tts.speak(f"Reminder: {msg}", lang)
        t = threading.Timer(minutes * 60, fire)
        t.daemon = True
        t.start()
        self._sched.append(t)
        unit = "minute" if minutes == 1 else "minutes"
        self.tts.speak(f"Reminder set. I will remind you about {msg} in {minutes} {unit}.", lang)

    def _start_scheduler(self):
        def run():
            while True:
                schedule.run_pending()
                time.sleep(10)
        threading.Thread(target=run, daemon=True).start()

    #  PowerPoint 
    def _create_ppt(self, p, lang):
        topic  = p.get("topic","Presentation")
        slides = int(p.get("slides", 5))
        self.tts.speak(f"Creating a presentation on {topic}.", lang)
        try:
            from pptx import Presentation as PPT
            from pptx.util import Pt
        except ImportError:
            self.tts.speak("Run: pip install python-pptx first.", lang)
            return

        prompt = (f"Create {slides} PowerPoint slides about: {topic}\n"
                  f"Return ONLY a JSON array like:\n"
                  f'[{{"title":"Slide title","content":"3 short bullet points"}}]\n'
                  f"No markdown. Just the JSON array.")
        raw = self.llm.simple_query(prompt)
        try:
            start = raw.find("[")
            end   = raw.rfind("]") + 1
            data  = json.loads(raw[start:end])
        except Exception:
            data = [
                {"title": topic,         "content": "Overview and introduction"},
                {"title": "Key Points",  "content": "Main concepts and ideas"},
                {"title": "Details",     "content": "In-depth information"},
                {"title": "Examples",    "content": "Real world applications"},
                {"title": "Summary",     "content": "Conclusion and next steps"},
            ]

        prs   = PPT()
        # Title slide
        sl    = prs.slides.add_slide(prs.slide_layouts[0])
        sl.shapes.title.text = topic
        try: sl.placeholders[1].text = f"Prepared by {APP_NAME}"
        except Exception: pass
        # Content slides
        for sd in data:
            sl  = prs.slides.add_slide(prs.slide_layouts[1])
            sl.shapes.title.text = sd.get("title","")
            try: sl.placeholders[1].text = sd.get("content","")
            except Exception: pass

        name  = topic.replace(" ","_").replace("/","_")[:30]
        out   = CODE_DIR / f"{name}.pptx"
        prs.save(str(out))
        try: os.startfile(str(out))
        except Exception: pass
        self.tts.speak(f"Presentation on {topic} is ready and opened.", lang)

    #  Research / Learn 
    def _research(self, p, lang):
        topic = p.get("topic","")
        self.tts.speak(f"Researching {topic} now. I will tell you when done.", lang)
        learner = self._get_learner()

        def on_progress(msg):
            # Only speak final done message
            if msg.startswith("Done.") or msg.startswith("Research complete"):
                self.tts.speak(msg, lang)

        learner.learn(topic, on_progress=on_progress, on_done=lambda m: None)

    #  Analyze file 
    def _analyze_file(self, p, lang):
        path = p.get("path","")
        fp   = Path(path)
        if not fp.exists():
            self.tts.speak(f"Cannot find {path}.", lang)
            return
        self.tts.speak("Reading the file.", lang)
        try:
            content = fp.read_text(encoding="utf-8", errors="replace")
            prompt  = (f"Analyze this code. In 3 sentences suggest improvements. "
                       f"Plain text only.\n\n{content[:3000]}")
            ans = self.llm.simple_query(prompt)
            self.tts.speak(ans, lang)
        except Exception as e:
            self.tts.speak(f"Could not read file: {str(e)[:60]}", lang)

    #  System info 
    def _sysinfo(self, p, lang):
        try:
            import psutil
            cpu  = psutil.cpu_percent(interval=1)
            ram  = psutil.virtual_memory().percent
            disk = psutil.disk_usage("/").percent
            self.tts.speak(f"CPU at {cpu:.0f}%, RAM at {ram:.0f}%, disk at {disk:.0f}%.", lang)
        except Exception:
            self.tts.speak("Could not read system metrics.", lang)

    #  PC control 
    def _restart(self, p, lang):
        self.tts.speak("Restarting in 5 seconds.", lang)
        time.sleep(5)
        os.system("shutdown /r /t 0" if sys.platform=="win32" else "sudo reboot")

    def _shutdown(self, p, lang):
        self.tts.speak("Shutting down now.", lang)
        time.sleep(3)
        os.system("shutdown /s /t 0" if sys.platform=="win32" else "sudo shutdown now")

    #  Teams 
    def _teams_meeting(self, p, lang):
        title = p.get("title","Meeting")
        start = p.get("start","")
        end   = p.get("end","")
        self.tts.speak(f"Creating Teams meeting: {title}.", lang)
        if self._win_integration():
            ok = self._teams.create_meeting(title, start, end)
            if ok:
                self.tts.speak("Teams meeting dialog opened. Fill in the details and send.", lang)
                return
        # Fallback: open Teams web
        q = urllib.parse.quote(title) if 'urllib' in sys.modules else title.replace(" ","+")
        webbrowser.open("https://teams.microsoft.com/l/meeting/new")
        self.tts.speak("Opened Teams new meeting page. Please fill in the details.", lang)

    #  Email 
    def _check_email(self, p, lang):
        from_addr = p.get("from","")
        self.tts.speak("Checking your email.", lang)
        if self._win_integration():
            result = self._outlook.check_latest_email(from_addr)
            if "error" in result:
                # Fallback: open Outlook or Gmail
                self.tts.speak("Opening Outlook for you.", lang)
                try:
                    subprocess.Popen(["outlook.exe"])
                except Exception:
                    webbrowser.open("https://mail.google.com")
            else:
                unread = "unread" if result.get("unread") else "read"
                msg = (f"Latest email from {result['from']}: "
                       f"Subject: {result['subject']}. "
                       f"Received {result['received'][:10]}. "
                       f"Message is {unread}.")
                self.tts.speak(msg, lang)
        else:
            self.tts.speak("Opening your email now.", lang)
            webbrowser.open("https://mail.google.com" if "@gmail" in from_addr
                            else "https://outlook.live.com")

    def _send_email(self, p, lang):
        to      = p.get("to","")
        subject = p.get("subject","")
        body    = p.get("body","")
        if self._win_integration():
            ok = self._outlook.send_email(to, subject, body)
            if ok:
                self.tts.speak(f"Email sent to {to}.", lang)
                return
        # Fallback: open mailto
        import urllib.parse
        url = f"mailto:{to}subject={urllib.parse.quote(subject)}&body={urllib.parse.quote(body)}"
        webbrowser.open(url)
        self.tts.speak(f"Opened email client. Please review and send.", lang)

    #  Screen analyze 
    def _screen_analyze(self, p, lang):
        self.tts.speak("Let me look at your screen.", lang)
        try:
            analyze = self._get_screen()
            result  = analyze(self.llm)
            self.tts.speak(result, lang)
        except Exception as e:
            self.tts.speak(f"Could not analyze screen: {str(e)[:60]}", lang)

    #  Model update 
    def _update_model(self, p, lang):
        self.tts.speak("Starting model update check. This runs in background and may take a while.", lang)
        def run():
            try:
                from core.model_updater import run as do_update
                do_update()
                flag = Path("data/model_updated.flag")
                if flag.exists():
                    new_model = flag.read_text().strip()
                    flag.unlink()
                    self.tts.speak(f"Model updated to {new_model}. Restart me to use it.", lang)
                else:
                    self.tts.speak("Already on the best available model.", lang)
            except Exception as e:
                self.tts.speak(f"Update check failed: {str(e)[:60]}", lang)
        threading.Thread(target=run, daemon=True).start()

    #  Dictation 
    def _dictate(self, p, lang):
        app = p.get("app","notepad")
        self.tts.speak(f"Opening {app} for dictation. Start speaking after the beep.", lang)
        try:
            import pyautogui
            app_map = {"notepad":["notepad.exe"],"wordpad":["wordpad.exe"]}
            subprocess.Popen(app_map.get(app.lower(),["notepad.exe"]))
            time.sleep(2)
            pyautogui.hotkey("alt","tab")
            time.sleep(0.5)
            self.tts.speak("Go. Say stop dictating when done.", lang)
        except Exception as e:
            self.tts.speak(f"Could not open {app}: {str(e)[:40]}", lang)

    # ==================== NEW COMMANDS ====================

    #  Identity response 
    def _identity(self, p, lang):
        # Get identity dynamically from config
        from core.config_generator import get_identity
        identity = get_identity()
        response = f"I am {identity['ai_name']}, personal assistant for {identity['user_name']}. I was created by {identity['creator']}."
        self.tts.speak(response, lang)

    #  Camera - take photo 
    def _take_photo(self, p, lang):
        self.tts.speak("Taking a photo.", lang)
        try:
            from core.media_recorder import get_camera_recorder
            camera = get_camera_recorder()
            filepath = camera.take_photo()
            if filepath:
                self.tts.speak(f"Photo saved to your Pictures folder.", lang)
            else:
                self.tts.speak("Could not take photo. Check if camera is connected.", lang)
        except Exception as e:
            self.tts.speak(f"Camera error: {str(e)[:40]}", lang)

    #  Camera - start video 
    def _start_video(self, p, lang):
        self.tts.speak("Starting video recording. Say stop recording when done.", lang)
        try:
            from core.media_recorder import get_camera_recorder
            camera = get_camera_recorder()
            if camera.start_video_recording():
                self.tts.speak("Recording.", lang)
            else:
                self.tts.speak("Could not start recording.", lang)
        except Exception as e:
            self.tts.speak(f"Error: {str(e)[:40]}", lang)

    #  Camera - stop video 
    def _stop_video(self, p, lang):
        try:
            from core.media_recorder import get_camera_recorder
            camera = get_camera_recorder()
            filepath = camera.stop_video_recording()
            if filepath:
                self.tts.speak(f"Video saved to your Videos folder.", lang)
            else:
                self.tts.speak("No recording in progress.", lang)
        except Exception as e:
            self.tts.speak(f"Error: {str(e)[:40]}", lang)

    #  Screenshot 
    def _screenshot(self, p, lang):
        self.tts.speak("Taking screenshot.", lang)
        try:
            from core.media_recorder import get_screen_recorder
            screen = get_screen_recorder()
            filepath = screen.take_screenshot()
            if filepath:
                self.tts.speak("Screenshot saved to your Pictures folder.", lang)
            else:
                self.tts.speak("Could not take screenshot.", lang)
        except Exception as e:
            self.tts.speak(f"Error: {str(e)[:40]}", lang)

    #  Screen recording - start 
    def _start_screen_rec(self, p, lang):
        self.tts.speak("Starting screen recording. Say stop screen recording when done.", lang)
        try:
            from core.media_recorder import get_screen_recorder
            screen = get_screen_recorder()
            if screen.start_screen_recording():
                self.tts.speak("Recording screen.", lang)
            else:
                self.tts.speak("Could not start screen recording.", lang)
        except Exception as e:
            self.tts.speak(f"Error: {str(e)[:40]}", lang)

    #  Screen recording - stop 
    def _stop_screen_rec(self, p, lang):
        try:
            from core.media_recorder import get_screen_recorder
            screen = get_screen_recorder()
            filepath = screen.stop_screen_recording()
            if filepath:
                self.tts.speak("Screen recording saved to your Videos folder.", lang)
            else:
                self.tts.speak("No screen recording in progress.", lang)
        except Exception as e:
            self.tts.speak(f"Error: {str(e)[:40]}", lang)

    #  Auto-learn (deep research and permanent storage) 
    def _auto_learn(self, p, lang):
        topic = p.get("topic", "")
        depth = p.get("depth", "normal")
        
        if not topic:
            self.tts.speak("What would you like me to learn about?", lang)
            return
        
        self.tts.speak(f"I will learn everything about {topic} and remember it forever. This may take a minute.", lang)
        
        learner = self._get_learner()
        
        def on_progress(msg):
            if "Done" in msg or "learned" in msg.lower():
                self.tts.speak(msg, lang)
        
        def on_complete(msg):
            self.tts.speak(f"I have learned about {topic}. You can ask me anything about it now, and I will remember forever.", lang)
        
        # Use enhanced learning (more sources, deeper storage)
        learner.learn(topic, on_progress=on_progress, on_done=on_complete)


import urllib.parse  # needed by teams/email fallbacks
